#!/usr/bin/env python3
"""
Generate PowerPoint Presentation for Mikroklima Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    return slide

def add_bullet(text_frame, text, level=0):
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(18) if level == 0 else Pt(16)
    return p

# Slide 1: Title
slide1 = add_title_slide(
    "Mikroklima Hamburg",
    "IoT Integration: FROST Server, ThingsBoard & InfluxDB\n\nAbdelrahman Ahmed (AAMHA)\nHochschule Osnabrück\n2026"
)

# Slide 2: Problem Statement
slide2 = add_content_slide("Research Questions")
content = slide2.placeholders[1].text_frame
content.clear()
add_bullet(content, "Can we validate sensor data quality?", 0)
add_bullet(content, "Compare citizen science sensors with professional data", 1)
add_bullet(content, "Result: 99.7% correlation (r=0.997) ✓", 1)
add_bullet(content, "", 0)
add_bullet(content, "Which IoT platform is best?", 0)
add_bullet(content, "Test 3 platforms with REAL data", 1)
add_bullet(content, "FROST Server (OGC), InfluxDB, ThingsBoard", 1)

# Slide 3: System Architecture
slide3 = add_content_slide("System Architecture")
content = slide3.placeholders[1].text_frame
content.clear()
add_bullet(content, "Data Sources → ETL Pipeline → 3 Platforms", 0)
add_bullet(content, "", 0)
add_bullet(content, "Real Data Sources (3):", 0)
add_bullet(content, "OpenSenseMap Hamburg (real-time PM10, PM2.5, Temp)", 1)
add_bullet(content, "Mobilithek Dormagen (11,139 measurements)", 1)
add_bullet(content, "Open-Meteo Egypt Cairo (192 weather records)", 1)
add_bullet(content, "", 0)
add_bullet(content, "Demo Sources (3): DWD, UDP Osnabrück, Tunisia", 0)

# Slide 4: Data Sources Detail
slide4 = add_content_slide("Real Data Sources")
content = slide4.placeholders[1].text_frame
content.clear()
add_bullet(content, "OpenSenseMap Hamburg", 0)
add_bullet(content, "Live API: PM10, PM2.5, Temperature, Humidity", 1)
add_bullet(content, "Citizen science project with real sensors", 1)
add_bullet(content, "", 0)
add_bullet(content, "Mobilithek Dormagen", 0)
add_bullet(content, "11,139 measurements from 8 sensors over 7 days", 1)
add_bullet(content, "Official German open data portal", 1)
add_bullet(content, "", 0)
add_bullet(content, "Open-Meteo Egypt (Cairo)", 0)
add_bullet(content, "192 professional weather measurements", 1)
add_bullet(content, "Used for ground truth validation", 1)

# Slide 5: Technical Implementation
slide5 = add_content_slide("Technical Implementation")
content = slide5.placeholders[1].text_frame
content.clear()
add_bullet(content, "Docker Architecture:", 0)
add_bullet(content, "7 containers: PostgreSQL, FROST, InfluxDB, ThingsBoard, Grafana", 1)
add_bullet(content, "", 0)
add_bullet(content, "Python ETL Pipeline:", 0)
add_bullet(content, "complete_data_loader.py - Fetches and pushes data", 1)
add_bullet(content, "activate_all_devices.py - Demo device activation", 1)
add_bullet(content, "temperature_comparison_egypt.py - Statistical analysis", 1)
add_bullet(content, "check_all_data.py - System verification", 1)
add_bullet(content, "", 0)
add_bullet(content, "Real-time processing: < 5 seconds per cycle", 0)

# Slide 6: Platform Comparison
slide6 = add_content_slide("IoT Platform Comparison")
content = slide6.placeholders[1].text_frame
content.clear()
add_bullet(content, "FROST Server (Port 8091)", 0)
add_bullet(content, "OGC SensorThings API standard (certified)", 1)
add_bullet(content, "Best for: Interoperability, international standards", 1)
add_bullet(content, "23 Things, 440 Observations", 1)
add_bullet(content, "", 0)
add_bullet(content, "InfluxDB (Port 8086)", 0)
add_bullet(content, "Time-series database with Flux query language", 1)
add_bullet(content, "Best for: Analytics, time-series operations", 1)
add_bullet(content, "8 sources, 97 measurements (24h)", 1)
add_bullet(content, "", 0)
add_bullet(content, "ThingsBoard (Port 8080)", 0)
add_bullet(content, "IoT platform with real-time dashboards", 1)
add_bullet(content, "Best for: Visualization, device management", 1)
add_bullet(content, "8 active devices with live telemetry", 1)

# Slide 7: Data Quality Analysis
slide7 = add_content_slide("Data Quality Analysis")
content = slide7.placeholders[1].text_frame
content.clear()
add_bullet(content, "Completeness: 85.7% ✓ GOOD", 0)
add_bullet(content, "11,139 total measurements", 1)
add_bullet(content, "1,568 missing values (14.3%)", 1)
add_bullet(content, "", 0)
add_bullet(content, "Quality Checks Performed:", 0)
add_bullet(content, "Gap detection and documentation", 1)
add_bullet(content, "Outlier identification", 1)
add_bullet(content, "Temporal consistency validation", 1)
add_bullet(content, "", 0)
add_bullet(content, "Conclusion: Data is reliable for analysis", 0)

# Slide 8: Temperature Validation (Egypt)
slide8 = add_content_slide("Temperature Validation - Egypt")
content = slide8.placeholders[1].text_frame
content.clear()
add_bullet(content, "Comparison: Citizen Sensors vs Professional Data", 0)
add_bullet(content, "", 0)
add_bullet(content, "Statistical Results:", 0)
add_bullet(content, "MAE: 0.75°C (Mean Absolute Error)", 1)
add_bullet(content, "RMSE: 0.83°C (Root Mean Square Error)", 1)
add_bullet(content, "Correlation: r = 0.997 (99.7%!) ⭐", 1)
add_bullet(content, "P-value: 0.00 (statistically significant)", 1)
add_bullet(content, "Sample size: 320 hourly measurements", 1)
add_bullet(content, "", 0)
add_bullet(content, "Conclusion: Citizen science sensors are RELIABLE!", 0)

# Slide 9: Results & Visualizations
slide9 = add_content_slide("Results & Visualizations")
content = slide9.placeholders[1].text_frame
content.clear()
add_bullet(content, "Generated Outputs:", 0)
add_bullet(content, "", 0)
add_bullet(content, "temperature_comparison_egypt.png", 0)
add_bullet(content, "3-panel chart: time series, difference, correlation", 1)
add_bullet(content, "", 0)
add_bullet(content, "sensor_locations_map.html", 0)
add_bullet(content, "Interactive map showing all sensor locations", 1)
add_bullet(content, "", 0)
add_bullet(content, "Live Dashboards:", 0)
add_bullet(content, "ThingsBoard: Real-time device monitoring", 1)
add_bullet(content, "InfluxDB: Time-series data explorer", 1)
add_bullet(content, "FROST: OGC-compliant API responses", 1)

# Slide 10: Key Achievements
slide10 = add_content_slide("Key Achievements")
content = slide10.placeholders[1].text_frame
content.clear()
add_bullet(content, "✓ Collected 11,331+ real measurements", 0)
add_bullet(content, "✓ Integrated 3 IoT platforms simultaneously", 0)
add_bullet(content, "✓ Validated sensor accuracy (99.7% correlation)", 0)
add_bullet(content, "✓ Followed OGC international standards", 0)
add_bullet(content, "✓ Dockerized production-ready system", 0)
add_bullet(content, "✓ Comprehensive statistical analysis", 0)
add_bullet(content, "✓ Real-time data pipeline (< 5 sec)", 0)
add_bullet(content, "✓ Full documentation and visualization", 0)

# Slide 11: Conclusion
slide11 = add_content_slide("Conclusion")
content = slide11.placeholders[1].text_frame
content.clear()
add_bullet(content, "Successfully compared 3 IoT platforms", 0)
add_bullet(content, "Each platform has specific strengths", 1)
add_bullet(content, "All work with real environmental data", 1)
add_bullet(content, "", 0)
add_bullet(content, "Proven sensor reliability", 0)
add_bullet(content, "Citizen science data quality validated", 1)
add_bullet(content, "r=0.997 correlation with professional sensors", 1)
add_bullet(content, "", 0)
add_bullet(content, "Production-ready system", 0)
add_bullet(content, "Docker containers, Python ETL, real-time processing", 1)
add_bullet(content, "GitHub: github.com/JUPA8/Mikroklima-Thingsboard-FROST", 1)

# Slide 12: Q&A
slide12 = add_title_slide(
    "Thank You!",
    "Questions?\n\n🌐 Live Demo Available\n📊 All platforms running\n📁 GitHub: JUPA8/Mikroklima-Thingsboard-FROST"
)

# Save presentation
prs.save('Mikroklima_Presentation.pptx')
print("✓ Presentation created: Mikroklima_Presentation.pptx")
print("\nSlides created:")
print("1. Title")
print("2. Research Questions")
print("3. System Architecture")
print("4. Real Data Sources")
print("5. Technical Implementation")
print("6. Platform Comparison")
print("7. Data Quality Analysis")
print("8. Temperature Validation")
print("9. Results & Visualizations")
print("10. Key Achievements")
print("11. Conclusion")
print("12. Q&A")
